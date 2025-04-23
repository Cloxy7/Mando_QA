import os
import textract
import pandas as pd
from PyPDF2 import PdfReader
import docx
import pptx
import json

def detect_and_parse(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif ext == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string()

        elif ext == ".pdf":
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

        elif ext == ".docx":
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == ".pptx":
            presentation = pptx.Presentation(file_path)
            return "\n".join(
                shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
            )

        elif ext == ".xlsx":
            df = pd.read_excel(file_path)
            return df.to_string()

        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
                return json.dumps(data, indent=2)

        else:
            # Fallback to textract
            return textract.process(file_path).decode("utf-8")

    except Exception as e:
        return f"❌ Error extracting text: {e}"
